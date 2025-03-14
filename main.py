import argparse

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN


def main(
    input_txt_file_path: Path,
    output_file_path: Path,
    font_size: float | int = 60,
    keep_blank_line: bool = False,
):
    # read txt line by line
    lyrics: list = input_txt_file_path.read_text(encoding="utf8")

    # delete blank lines
    lyrics = lyrics.split("\n")
    if not keep_blank_line:
        lyrics = [line for line in lyrics if line != ""]

    # create ppt file
    root = Presentation()
    # 16:9
    root.slide_width = Inches(13.33)
    root.slide_height = Inches(7.5)

    # blank layout
    slide_layout = root.slide_layouts[6]

    # write slides line by line
    for line in lyrics:
        slide = root.slides.add_slide(slide_layout)

        # set background to black
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 0, 0)

        # creating textFrames
        left = top = Inches(0.25)
        width = root.slide_width - Inches(0.5)
        height = root.slide_height - Inches(0.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.add_paragraph()

        # add the line of lyrics to the paragraph
        p.text = line

        # set font and color
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True

        # tf.fit_text()

    # output the ppt file
    root.save(output_file_path)
    print(f"Successfully outputted to {output_file_path}")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(
        description="Convert lyrics text file to PowerPoint presentation. Outputs a file with the same name as the input file by default."
    )
    parser.add_argument(
        "input_txt_file_path", type=Path, help="Path to the input lyrics text file."
    )
    parser.add_argument(
        "--font_size", type=float, default=60, help="Font size for the lyrics text, unit in Pt, defaults to 60."
    )
    parser.add_argument(
        "--preserve_newline",
        action="store_true",
        help="Preserve blank lines in the input txt file. Defaults to skip blank lines.",
    )

    args = parser.parse_args()

    input_txt_file_path = Path(args.input_txt_file_path)
    if not input_txt_file_path.exists():
        print(f"The specified input file {input_txt_file_path} does not exist.")

    output_file_path = input_txt_file_path.with_suffix(".pptx")
    if output_file_path.exists():
        print(
            f"The output file {input_txt_file_path} already exists. Will be overwriting the file!"
        )

    main(
        args.input_txt_file_path,
        output_file_path,
        args.font_size,
        args.preserve_newline,
    )
